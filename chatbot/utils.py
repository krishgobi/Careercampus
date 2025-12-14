"""
Utility functions for RAG (Retrieval-Augmented Generation) system
"""
import os
import pickle
from pathlib import Path
from typing import List, Tuple
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
import numpy as np
from django.conf import settings

# API clients - initialized lazily to avoid import-time errors
_groq_client = None
_vectorizer = None
_vectorizer_fitted = False

# Vector store directory - not needed for TF-IDF but keeping for compatibility
VECTOR_STORE_DIR = Path(settings.MEDIA_ROOT) / 'vector_stores'
VECTOR_STORE_DIR.mkdir(exist_ok=True, parents=True)


def call_llm_api(model_id, messages, temperature=0.7, max_tokens=1024):
    """
    Call LLM API for chat completions - supports both Gemini and Groq providers
    
    Args:
        model_id: Model identifier (e.g., 'gemini-2.5-flash', 'llama-3.1-8b-instant')
        messages: List of message dicts with 'role' and 'content'
        temperature: Temperature for generation
        max_tokens: Maximum tokens to generate
    
    Returns:
        Generated text response
    """
    # Determine provider based on model_id
    gemini_models = ['gemini-2.5-flash', 'gemini-2.0-flash-exp', 'gemini-2.0-flash-lite']
    
    if model_id in gemini_models:
        return _call_gemini(model_id, messages, temperature, max_tokens)
    else:
        return _call_groq(model_id, messages, temperature, max_tokens)


def _call_gemini(model_id, messages, temperature, max_tokens):
    """Call Gemini API"""
    import google.generativeai as genai
    
    # Configure Gemini
    genai.configure(api_key=settings.GEMINI_API_KEY)
    
    # Use Gemini 2.5 Flash if the default model is specified
    if model_id in ['gemini-2.0-flash-exp', 'gemini-2.0-flash-lite']:
        model_id = 'gemini-2.5-flash'
    
    # Convert messages to Gemini format
    prompt_parts = []
    for msg in messages:
        role = msg.get("role", "")
        content = msg.get("content", "")
        if role == "system":
            prompt_parts.append(f"Instructions: {content}\n")
        elif role == "user":
            prompt_parts.append(f"User: {content}\n")
        elif role == "assistant":
            prompt_parts.append(f"Assistant: {content}\n")
    
    prompt = "\n".join(prompt_parts)
    
    # Create model instance
    model = genai.GenerativeModel(model_id)
    
    response = model.generate_content(
        prompt,
        generation_config=genai.types.GenerationConfig(
            temperature=temperature,
            max_output_tokens=max_tokens,
        )
    )
    
    return response.text


def _call_groq(model_id, messages, temperature, max_tokens):
    """Call Groq API"""
    import requests
    
    url = "https://api.groq.com/openai/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {settings.GROQ_API_KEY}",
        "Content-Type": "application/json"
    }
    
    payload = {
        "model": model_id,
        "messages": messages,
        "temperature": temperature,
        "max_tokens": max_tokens
    }
    
    response = requests.post(url, headers=headers, json=payload)
    
    if response.status_code != 200:
        print(f"[ERROR] Groq API returned {response.status_code}")
        print(f"[ERROR] Response: {response.text}")
        raise Exception(f"Groq API error: {response.text}")
    
    return response.json()["choices"][0]["message"]["content"]


def get_vectorizer():
    """Get TF-IDF vectorizer for embeddings"""
    global _vectorizer
    if _vectorizer is None:
        from sklearn.feature_extraction.text import TfidfVectorizer
        print("[INFO] Initializing TF-IDF vectorizer...")
        _vectorizer = TfidfVectorizer(
            max_features=384,  # Keep dimensionality similar to transformer models
            ngram_range=(1, 2),  # Unigrams and bigrams
            min_df=1,  # Minimum document frequency
            max_df=1.0  # Maximum document frequency
        )
        print("[INFO] Vectorizer ready!")
    return _vectorizer



def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file"""
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        print(f"Error extracting PDF: {e}")
    return text


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file"""
    text = ""
    try:
        doc = DocxDocument(file_path)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    except Exception as e:
        print(f"Error extracting DOCX: {e}")
    return text


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file"""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error extracting PPTX: {e}")
    return text


def extract_text_from_file(file_path: str, file_type: str) -> str:
    """Extract text based on file type"""
    extractors = {
        'pdf': extract_text_from_pdf,
        'docx': extract_text_from_docx,
        'pptx': extract_text_from_pptx,
    }
    
    extractor = extractors.get(file_type.lower())
    if extractor:
        return extractor(file_path)
    return ""


def chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
    """Split text into chunks for better retrieval"""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    chunks = text_splitter.split_text(text)
    return chunks


def get_embeddings(texts: List[str]) -> np.ndarray:
    """Generate embeddings using TF-IDF (fast, simple, no API/quota needed)"""
    vectorizer = get_vectorizer()
    
    print(f"[DEBUG] Generating TF-IDF embeddings for {len(texts)} chunks...")
    
    try:
        # Fit and transform the texts
        embeddings = vectorizer.fit_transform(texts).toarray()
        print(f"[DEBUG] Generated embeddings with shape: {embeddings.shape}")
        return embeddings.astype('float32')
    except Exception as e:
        print(f"[ERROR] Error generating embeddings: {type(e).__name__}: {str(e)}")
        # Return zero vectors as fallback
        return np.zeros((len(texts), 384), dtype='float32')


def create_vector_store(document_id: int, chunks: List[str]) -> Tuple[None, List[str]]:
    """Create vector store for document chunks (saves fitted vectorizer for consistent dimensions)"""
    if not chunks:
        return None, []
    
    # Create and fit vectorizer on document chunks
    from sklearn.feature_extraction.text import TfidfVectorizer
    vectorizer = TfidfVectorizer(
        max_features=384,
        ngram_range=(1, 2),
        min_df=1,
        max_df=1.0
    )
    
    # Fit and transform the chunks
    print(f"[DEBUG] Fitting TF-IDF vectorizer on {len(chunks)} chunks...")
    embeddings = vectorizer.fit_transform(chunks).toarray()
    print(f"[DEBUG] Created embeddings with shape: {embeddings.shape}")
    
    # Save embeddings, chunks, AND the fitted vectorizer
    embeddings_path = VECTOR_STORE_DIR / f'doc_{document_id}_embeddings.pkl'
    chunks_path = VECTOR_STORE_DIR / f'doc_{document_id}_chunks.pkl'
    vectorizer_path = VECTOR_STORE_DIR / f'doc_{document_id}_vectorizer.pkl'
    
    with open(embeddings_path, 'wb') as f:
        pickle.dump(embeddings, f)
    with open(chunks_path, 'wb') as f:
        pickle.dump(chunks, f)
    with open(vectorizer_path, 'wb') as f:
        pickle.dump(vectorizer, f)
    
    print(f"[INFO] Saved {len(chunks)} chunks, embeddings, and vectorizer for document {document_id}")
    return None, chunks


def load_vector_store(document_id: int) -> Tuple[np.ndarray, List[str], object]:
    """Load vector store and fitted vectorizer from disk"""
    embeddings_path = VECTOR_STORE_DIR / f'doc_{document_id}_embeddings.pkl'
    chunks_path = VECTOR_STORE_DIR / f'doc_{document_id}_chunks.pkl'
    vectorizer_path = VECTOR_STORE_DIR / f'doc_{document_id}_vectorizer.pkl'
    
    if not embeddings_path.exists() or not chunks_path.exists() or not vectorizer_path.exists():
        return None, [], None
    
    with open(embeddings_path, 'rb') as f:
        embeddings = pickle.load(f)
    with open(chunks_path, 'rb') as f:
        chunks = pickle.load(f)
    with open(vectorizer_path, 'rb') as f:
        vectorizer = pickle.load(f)
    
    return embeddings, chunks, vectorizer


def retrieve_relevant_chunks(query: str, document_id: int, k: int = 3) -> List[str]:
    """Retrieve most relevant chunks for a query using cosine similarity"""
    from sklearn.metrics.pairwise import cosine_similarity
    
    embeddings, chunks, vectorizer = load_vector_store(document_id)
    
    if embeddings is None or not chunks or vectorizer is None:
        return []
    
    # Use the SAME fitted vectorizer to transform the query (ensures same dimensions)
    print(f"[DEBUG] Transforming query using saved vectorizer...")
    query_embedding = vectorizer.transform([query]).toarray()
    print(f"[DEBUG] Query embedding shape: {query_embedding.shape}, Document embeddings shape: {embeddings.shape}")
    
    # Calculate cosine similarity
    similarities = cosine_similarity(query_embedding, embeddings)[0]
    
    # Get top k indices
    top_indices = np.argsort(similarities)[::-1][:k]
    
    # Return relevant chunks
    relevant_chunks = [chunks[i] for i in top_indices if i < len(chunks)]
    return relevant_chunks


def generate_answer(query: str, context: str, model_id: str = 'llama-3.1-8b-instant', chat_history: List[dict] = None) -> str:
    """Generate answer using LLM with RAG context - supports multiple models with fallback"""
    if chat_history is None:
        chat_history = []
    
    # Prepare system message with context
    system_message = f"""You are a helpful AI assistant for students. You answer questions based on the provided context from their course materials.

Context from documents:
{context}

Instructions:
- Answer the question using ONLY the information from the context above
- If the context doesn't contain enough information, say "I don't have enough information in the uploaded documents to answer this question"
- Be clear, concise, and educational
- If relevant, provide examples or explanations to help the student understand better
"""
    
    # Build messages for API
    messages = [{"role": "system", "content": system_message}]
    
    # Add chat history (last 5 messages for context)
    for msg in chat_history[-5:]:
        messages.append({
            "role": msg["role"],
            "content": msg["content"]
        })
    
    # Add current query
    messages.append({"role": "user", "content": query})
    
    try:
        print(f"[DEBUG] Calling LLM API ({model_id}) with query: {query[:50]}...")
        # Call unified LLM API (supports Gemini + Groq)
        answer = call_llm_api(model_id, messages, temperature=0.7, max_tokens=1024)
        print(f"[DEBUG] LLM API response received: {answer[:100]}...")
        return answer
    
    except Exception as e:
        error_str = str(e).lower()
        print(f"[ERROR] Error generating answer: {type(e).__name__}: {str(e)}")
        
        # If quota exhausted on Gemini, try Groq
        if 'resource' in error_str or 'quota' in error_str or 'exhausted' in error_str:
            print("[INFO] Gemini quota exhausted, switching to Groq...")
            try:
                # Try with Groq model
                fallback_model = 'llama-3.3-70b-versatile'
                print(f"[INFO] Retrying with Groq model: {fallback_model}")
                answer = call_llm_api(fallback_model, messages, temperature=0.7, max_tokens=1024)
                return answer
            except Exception as fallback_error:
                print(f"[ERROR] Groq fallback also failed: {fallback_error}")
                return f"I'm sorry, both AI services are currently unavailable. Please try again later."
        
        # For other errors, return generic error
        import traceback
        traceback.print_exc()
        return f"I'm sorry, I encountered an error while generating the answer: {type(e).__name__}. Please check the server logs for details."


def process_query(query: str, document_id: int, chat_history: List[dict] = None) -> str:
    """Process a query using RAG pipeline"""
    print(f"[DEBUG] Processing query for document {document_id}: {query[:100]}...")
    
    # Retrieve relevant chunks
    relevant_chunks = retrieve_relevant_chunks(query, document_id, k=3)
    
    print(f"[DEBUG] Retrieved {len(relevant_chunks)} chunks")
    if relevant_chunks:
        for i, chunk in enumerate(relevant_chunks):
            print(f"[DEBUG] Chunk {i+1}: {chunk[:100]}...")
    
    if not relevant_chunks:
        print("[ERROR] No relevant chunks found!")
        return "I couldn't find relevant information in the document. Please make sure the document has been properly uploaded and processed."
    
    # Combine chunks into context
    context = "\n\n".join(relevant_chunks)
    print(f"[DEBUG] Context combined, total length: {len(context)} characters")
    
    # Generate answer
    answer = generate_answer(query, context, chat_history)
    
    return answer
